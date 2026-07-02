"""AI Service — Text Extractor (PDF, PPTX, TXT)."""

import io
import structlog

logger = structlog.get_logger()


class TextExtractor:
    """Extracts text from various file formats."""

    def extract(self, file_bytes: bytes, content_type: str, filename: str) -> str:
        """
        Extract text from file bytes based on content type.

        Args:
            file_bytes: Raw file content.
            content_type: MIME type of the file.
            filename: Original filename for fallback detection.

        Returns:
            Extracted text content.
        """
        if content_type == "application/pdf" or filename.lower().endswith(".pdf"):
            return self._extract_pdf(file_bytes)
        elif (
            content_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            or filename.lower().endswith(".pptx")
        ):
            return self._extract_pptx(file_bytes)
        elif content_type == "text/plain" or filename.lower().endswith(".txt"):
            return self._extract_text(file_bytes)
        else:
            raise ValueError(f"Unsupported file type: {content_type} ({filename})")

    def _extract_pdf(self, file_bytes: bytes) -> str:
        """Extract text from PDF using pdfplumber (fallback: PyPDF2)."""
        text = ""

        # Try pdfplumber first (better quality)
        try:
            import pdfplumber
            with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n\n"
            if text.strip():
                logger.info("pdf_extracted", method="pdfplumber", length=len(text))
                return text.strip()
        except Exception as e:
            logger.warning("pdfplumber_failed", error=str(e))

        # Fallback to PyPDF2
        try:
            from PyPDF2 import PdfReader
            reader = PdfReader(io.BytesIO(file_bytes))
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n\n"
            logger.info("pdf_extracted", method="PyPDF2", length=len(text))
            return text.strip()
        except Exception as e:
            logger.error("pdf_extraction_failed", error=str(e))
            raise ValueError(f"Failed to extract text from PDF: {e}")

    def _extract_pptx(self, file_bytes: bytes) -> str:
        """Extract text from PowerPoint PPTX."""
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(file_bytes))
            text_parts = []

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = f"--- Slide {slide_num} ---\n"
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            para_text = paragraph.text.strip()
                            if para_text:
                                slide_text += para_text + "\n"
                text_parts.append(slide_text)

            text = "\n".join(text_parts)
            logger.info("pptx_extracted", slides=len(prs.slides), length=len(text))
            return text.strip()
        except Exception as e:
            logger.error("pptx_extraction_failed", error=str(e))
            raise ValueError(f"Failed to extract text from PPTX: {e}")

    def _extract_text(self, file_bytes: bytes) -> str:
        """Extract text from plain text file."""
        try:
            text = file_bytes.decode("utf-8")
        except UnicodeDecodeError:
            text = file_bytes.decode("latin-1")

        logger.info("text_extracted", length=len(text))
        return text.strip()
